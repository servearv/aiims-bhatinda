import sys
try:
    from pptx import Presentation
except ImportError:
    import subprocess
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation

def read_ppt(file_path):
    try:
        ppt = Presentation(file_path)
        print(f"Successfully opened {file_path}. It has {len(ppt.slides)} slides.\n")
        
        for i, slide in enumerate(ppt.slides):
            print(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        print(text)
            print()
            
    except Exception as e:
        print(f"Error reading PPT: {e}")

if __name__ == "__main__":
    read_ppt('MedDigitizer.pptx')
